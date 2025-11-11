import os
import PyPDF2
from docx import Document
from pptx import Presentation
from typing import List, Dict
import re

class TextProcessor:
    """Process various text file formats"""
    
    @staticmethod
    def process_pdf(file_path: str) -> List[Dict]:
        """Extract text from PDF file"""
        chunks = []
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                full_text = ""
                
                for page_num, page in enumerate(pdf_reader.pages):
                    text = page.extract_text()
                    full_text += f"\n--- Page {page_num + 1} ---\n{text}\n"
                
                # Split into chunks
                chunks = TextProcessor._split_text(full_text)
                
        except Exception as e:
            print(f"Error processing PDF: {e}")
            chunks = [{"content": f"Error reading PDF: {str(e)}", "metadata": {}}]
        
        return chunks
    
    @staticmethod
    def process_docx(file_path: str) -> List[Dict]:
        """Extract text from DOCX file"""
        chunks = []
        try:
            doc = Document(file_path)
            full_text = ""
            
            for paragraph in doc.paragraphs:
                full_text += paragraph.text + "\n"
            
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        full_text += cell.text + " "
                full_text += "\n"
            
            chunks = TextProcessor._split_text(full_text)
            
        except Exception as e:
            print(f"Error processing DOCX: {e}")
            chunks = [{"content": f"Error reading DOCX: {str(e)}", "metadata": {}}]
        
        return chunks
    
    @staticmethod
    def process_pptx(file_path: str) -> List[Dict]:
        """Extract text from PPTX file"""
        chunks = []
        try:
            prs = Presentation(file_path)
            full_text = ""
            
            for slide_num, slide in enumerate(prs.slides):
                full_text += f"\n--- Slide {slide_num + 1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text += shape.text + "\n"
            
            chunks = TextProcessor._split_text(full_text)
            
        except Exception as e:
            print(f"Error processing PPTX: {e}")
            chunks = [{"content": f"Error reading PPTX: {str(e)}", "metadata": {}}]
        
        return chunks
    
    @staticmethod
    def process_txt(file_path: str) -> List[Dict]:
        """Extract text from TXT file"""
        chunks = []
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                content = file.read()
                chunks = TextProcessor._split_text(content)
        except Exception as e:
            print(f"Error processing TXT: {e}")
            chunks = [{"content": f"Error reading TXT: {str(e)}", "metadata": {}}]
        
        return chunks
    
    @staticmethod
    def process_md(file_path: str) -> List[Dict]:
        """Extract text from Markdown file"""
        chunks = []
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                content = file.read()
                chunks = TextProcessor._split_text(content)
        except Exception as e:
            print(f"Error processing MD: {e}")
            chunks = [{"content": f"Error reading MD: {str(e)}", "metadata": {}}]
        
        return chunks
    
    @staticmethod
    def _split_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> List[Dict]:
        """Split text into chunks with overlap"""
        if not text or not text.strip():
            return [{"content": "", "metadata": {}}]
        
        # Clean text
        text = re.sub(r'\s+', ' ', text).strip()
        
        chunks = []
        start = 0
        chunk_index = 0
        
        while start < len(text):
            end = start + chunk_size
            
            # Try to break at sentence boundary
            if end < len(text):
                # Look for sentence endings
                sentence_end = max(
                    text.rfind('.', start, end),
                    text.rfind('!', start, end),
                    text.rfind('?', start, end),
                    text.rfind('\n', start, end)
                )
                
                if sentence_end > start:
                    end = sentence_end + 1
            
            chunk_text = text[start:end].strip()
            if chunk_text:
                chunks.append({
                    "content": chunk_text,
                    "metadata": {
                        "chunk_index": chunk_index,
                        "start_pos": start,
                        "end_pos": end
                    }
                })
                chunk_index += 1
            
            # Move start position with overlap
            start = max(start + 1, end - overlap)
        
        return chunks if chunks else [{"content": text, "metadata": {}}]
    
    @staticmethod
    def process_file(file_path: str) -> List[Dict]:
        """Process a text file based on its extension"""
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext == '.pdf':
            return TextProcessor.process_pdf(file_path)
        elif ext == '.docx':
            return TextProcessor.process_docx(file_path)
        elif ext == '.pptx':
            return TextProcessor.process_pptx(file_path)
        elif ext == '.txt':
            return TextProcessor.process_txt(file_path)
        elif ext == '.md':
            return TextProcessor.process_md(file_path)
        else:
            return [{"content": f"Unsupported file type: {ext}", "metadata": {}}]

